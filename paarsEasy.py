import wx
import wx.html as wxhtml
import os
import pandas as pd
import numpy as np
import paramiko
import threading
import time
import datetime
import shutil
from openpyxl import load_workbook
import getpass
import json
import glob
from collections import OrderedDict, namedtuple
import pickle
from wx.lib import masked
from scipy import stats
from random import random
from tabulate import tabulate
from corr import standarderr
import dataframe_image as dfi
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
from shutil import copyfile

WILDCARD = "Excel sheets (*.xls;*.xlsx)|*.xls;*.xlsx"     
IPADDRESS = "azmlg01.southcentralus.cloudapp.azure.com"  #"ec2-3-22-226-74.us-east-2.compute.amazonaws.com"
HOMEDIR = '/home/working/' # NB home directory on Linux 
RESULTSDIR = HOMEDIR + 'results/' # NB home directory on Linux - these 2 not the Windows directories
#INSTANCE_ID = 'i-03889c587823b7ec5' 
DEBUG = False 

class constants(object):
    pass
K = constants()
K.csvfiles=[]

if getpass.getuser() == 'pellis': # then we are on Azure
    PCDIRECTORY = 'c:/PAARS/CDADesktop/'
    ROLLFORWARD = 'c:/PAARS/CDADesktop/DummyOutputSS.xlsx'
    KEYFILE = 'c:/PAARS/CDADesktop/azmlg01_key.pem'
    K.INDIR = 'c:/PAARS/' # the place where the raw economic scenario CSV files are store
else:
    print('Yikes - what happened to the Azure login.  Whoareyou?')
    #PCDIRECTORY = 'c:/PAARS/'
    #ROLLFORWARD = 'c:/PAARS/DummyOutputSS.xlsx'
    #KEYFILE = 'c:/PAARS/azmlg01_key.pem'
    #K.INDIR = 'c:/PAARS/'

KEY = paramiko.RSAKey.from_private_key_file(KEYFILE)
CLIENT = paramiko.SSHClient()
CLIENT.set_missing_host_key_policy(paramiko.AutoAddPolicy())

def flt(s):
    try:
        result = float(s)
    except:
        result = s
    return result 

def results_pic():
    #  create results table picture for putting at the end of the powerpoint file
    def sty(colname):
        if '$' in colname:
            return "${:,.0f}"
        if '%' in colname:
            return "{:.2f}%"
        if "year" in colname.lower():
            return "{:.0f}"
        if "number" in colname.lower():
            return "{:,.0f}"
        return "{}"
    t1 = 'ResultsTable.xlsx'
    #try: should not be possible can only get here after creating resultstable.xlsx
    df = pd.read_excel(t1,index_col=0)
    #except:
    #    return 'Fail'
    wanted_cols=['CTE 0 NPV CF $','CTE 70 $','CTE 98 $','C3P2 $','800RBC %','Worst NPV CF $']
    df = df[wanted_cols]
    styledict=dict((i,sty(i)) for i in df.columns)
    df2 = df.style.format(styledict)
    dfi.export(df2,'resultstable.png')
    return # used to return 'ok' if it did not fail above.  Nn now

def histo():
    # create 2 histograms with the spectrum of cashflows for putting at the end of the powerpoint
    pickledick = {}
    for fil in glob.glob('spectrumthisy' + '*.p'):
        pieces = fil[:-2].split('_') # remove the .p the last two chars
        nm = pieces[1]
        if len(pieces) > 2:
            sex = pieces[2][0]
            age = pieces[2][1:3]
            lis = [nm,age,sex]
        else:
            lis = [nm]
        name = '_'.join(lis)
        pickledick[name] = pickle.load(open(fil,"rb"))
    graphnames = []
    for k,v in pickledick.items():
        fig, ax = plt.subplots(figsize=(20,20))
        sns.displot(v/1000,kde=False,bins=40)
        titstr = 'Distribution of NPV Cashflows for '+k+' Fund ($000)'
        plt.title(titstr,fontsize=6)
        #ax.figure.tight_layout()
        filename = k + 'CashflowHistogram.png'
        graphnames.append(filename)
        plt.savefig(filename)
    graphnames.sort() # get 1 first and 2 next, for layout in orderly fashion
    return graphnames

def outputppt(destination):
    # add the monte carlo summary analysis to the end of the powerpoint file
    # Get the powerpoint file - whatever the first one is!
    try:
        pptfile =  glob.glob(destination + '*.pptx')[0]
    except:
        pptfile = destination + 'montecarlo.ppt'
        copyfile('c:/users/pellis/desktop/blankpres.ppt', pptfile)
    prs = Presentation(pptfile)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    top,left = Inches(0.5), Inches(2)
    slide.shapes.add_picture('c:/PAARS/CDADesktop/MonteCarlo.png', left, top)
    top, left = Inches(2), Inches(4)
    slide.shapes.add_picture('resultstable.png', left, top,height=Inches(1))
    topinch, leftinch = 4, 1
    for gn in graphnames:
        top = Inches(topinch)
        left = Inches(leftinch)
        pic = slide.shapes.add_picture(gn, left, top, height=Inches(2.5),width=Inches(4))
        leftinch += 5
    prs.save(pptfile)
    return(pptfile)

class MyHtmlFrame(wx.Dialog):
    # Simple HTML message
    def __init__(self,parent,title,content):
        wx.Frame.__init__(self,parent,-1,title)
        htm = wxhtml.HtmlWindow(self) 
        if "gtk2" in wx.PlatformInfo: htm.SetStandardFonts()
        htm.SetPage(content)
        self.SetSize((500,500))
        self.Fit()

def htmlmsg(title,content):
    frm = MyHtmlFrame(None, title,content)
    frm.ShowModal()
    frm.Destroy()

def grow(df): # extract growth rates from cumulative values
    # NB the first column is a zero - this is compliant with the non-cumulative Conning files
    #  And the way the algorithm works, this will give you a 1 as the first year value
    #  when turned back to cumulative.
    num = df.iloc[:,1:]
    den = df.iloc[:,:-1]
    den.columns = num.columns
    growth = num/den - 1
    growth.insert(0,0,0) # insert column zero at position zero with zeros in it!
    return growth

def undermusig(ez, sdz): # gets the underlying mu and sigma for the lognormal
    ez2 = ez * ez
    varz = sdz * sdz
    ex = np.log(ez2/np.sqrt(ez2 + varz))
    sdx = np.sqrt(np.log(1+varz/ez2))
    return ex,sdx

def logn(p,mu,sigma):
    return np.log(stats.lognorm.ppf(p, sigma, loc=0, scale=np.exp(mu))) + 1

def loader(fn, indir,filetype):
    #print(fn)
    dfin = pd.read_csv(indir + fn + '.csv',header=None)
    if filetype == 'Conning':
        return dfin
    else:
        return grow(dfin) # gets the growth rates each period from the cumulative AAA file

def blendo(parent, data, indir, outdir, fname, error, pftype): # default error is 5%
    # files is a list.  [[asset name, weight]]
    # error is the error term as a fraction, derived from
    parent.frame.SetStatusText('Blending files - please wait.')
    filenms = [i[0] for i in data[:-2]] # last two items are total and r-squared
    pcts = [i[1] for i in data[:-2]]
    returns = [loader(fn,indir,pftype) for fn in filenms]
    balret = [ret * pct for ret,pct in zip(returns,pcts)]
    sumbalret = sum(balret)
    sds = error * abs(sumbalret)
    e1 = np.random.normal(sumbalret,sds)
    e2 = pd.DataFrame(e1)
    errs = round(e2,6)
    errs.columns = sumbalret.columns
    balgro = 1 + sumbalret
    errbalgro = 1 + errs
    balfundval = np.cumproduct(balgro,axis=1)
    errbalfundval = np.cumproduct(errbalgro,axis=1)
    cols_needed = [i for i in balfundval.columns if i % 12 == 0]
    outputbalfund=round(balfundval[cols_needed],6)
    #outputbalfund.insert(0,-1,1)
    #outputbalfund.to_csv(outdir + '//' + fname +'raw.csv', header=False, index=False)
    outputerrbal = round(errbalfundval[cols_needed],6)
    #outputerrbal.insert(0,-1,1)
    outputerrbal.to_csv(outdir + fname, header = False, index = False)
    #round(sumbalret,6).to_csv(outdir + '//' + fname + 'returns.csv', header = False, index = False)
    #round(errs,6).to_csv(outdir + '//' + fname + 'returnserr.csv', header = False, index = False)


class SummaryResult(wx.Dialog): 
    def __init__(self, parent): 
        super(SummaryResult, self).__init__(parent, title = 'Summary Results', size = (1100,300)) 
        panel = wx.Panel(self) 
		# get and format the results
        self.K = parent.K
        self.rslt = self.getresults()
        self.formatresults()
        self.btn = wx.Button(panel, wx.ID_OK, label = "ok") #, size = (50,50)) #, pos = (75,50))
        txt_style = wx.VSCROLL|wx.HSCROLL|wx.TE_READONLY|wx.BORDER_SIMPLE
        htwin = wxhtml.HtmlWindow(panel, -1, size=(1200, 1200), style=txt_style)
        htwin.SetPage(self.rslt.to_html())
        sizer = wx.BoxSizer(wx.VERTICAL)
        sizer.Add(htwin, 0, wx.ALL, 10)
        sizer.Add(self.btn, 0, wx.ALL, 10)
        panel.SetSizer(sizer)
        panel.Layout()
		
    def getresults(self):
        scenarios = json.load(open(self.K.outdir + 'scenario.json'))
        names = scenarios.keys()
        descs = [v['Description'] for v in scenarios.values()]
        resultsdata = OrderedDict()
        for n,d in zip(names,descs):
            fn = self.K.outdir + 'table_' + n + '.P'
            results = (pickle.load(open(fn,'rb')))
            results['Description'] = d
            resultsdata[n] = results
        resultsdf = pd.DataFrame(resultsdata)
        ix = list(resultsdf.index)
        # remove Description from wherever it is
        descrix = ix.index('Description')
        ix.pop(descrix)
        ix.insert(0,'Description') # is there a better way???
        resultsdf_ = resultsdf.reindex(ix) # now description is at top
        dft = resultsdf_.T
        dft.to_excel(self.K.outdir + 'ResultsTable.xlsx')
        return dft

    def formatresults(self):
        for k in self.rslt.columns:
            if '%' in k:
                fstr = "{:.2f}%"
            elif '$' in k:
                fstr = "${:,.0f}"
            elif 'year' in k.lower():
                fstr = "{:.0f}"
            elif 'number' in k.lower():
                fstr = "{:,.0f}"
            else:
                fstr = "{0}"
            self.rslt[k] = [fstr.format(i) for i in self.rslt[k]]
"""

class SummaryResult(object): 
    def __init__(self, parent): 
        self.K = parent.K
        self.rslt = self.getresults()
        self.formatresults()
        htmlmsg('Results Table',self.rslt.to_html())
		
    def getresults(self):
        scenarios = json.load(open(self.K.outdir + 'scenario.json'))
        names = scenarios.keys()
        descs = [v['Description'] for v in scenarios.values()]
        resultsdata = OrderedDict()
        for n,d in zip(names,descs):
            fn = self.K.outdir + 'table_' + n + '.P'
            results = (pickle.load(open(fn,'rb')))
            results['Description'] = d
            resultsdata[n] = results
        resultsdf = pd.DataFrame(resultsdata)
        ix = list(resultsdf.index)
        # remove Description from wherever it is
        descrix = ix.index('Description')
        ix.pop(descrix)
        ix.insert(0,'Description') # is there a better way???
        resultsdf_ = resultsdf.reindex(ix) # now description is at top
        dft = resultsdf_.T
        dft.to_excel(self.K.outdir + 'ResultsTable.xlsx')
        return dft

    def formatresults(self):
        for k in self.rslt.columns:
            if '%' in k:
                fstr = "{:.2f}%"
            elif '$' in k:
                fstr = "${:,.0f}"
            elif 'year' in k.lower():
                fstr = "{:.0f}"
            elif 'number' in k.lower():
                fstr = "{:,.0f}"
            else:
                fstr = "{0}"
            self.rslt[k] = [fstr.format(i) for i in self.rslt[k]]
"""

class Monty(object): # Start the Monte Carlo process

    def __init__(self,parent,scenarios,K):
        self.parent = parent
        self.scenarios = scenarios # scenarios object
        self.K = K
        self.sendfile()
        if self.msg == '':
            t = threading.Thread(target=self.worker)
            u = threading.Thread(target=self.checker)
            t.start()
            u.start()
        else:
            self.parent.frame.msg(self.msg)

    def connectme(self):
        CLIENT.connect(hostname=IPADDRESS, username="meritmodeling", pkey=KEY)

    def commandme(self,cmd):
        stdin, stdout, stderr = CLIENT.exec_command(cmd)
        if DEBUG:
            print('cmd',cmd)
            self.pri('stdout',stdout)
            self.pri('stderr',stderr)

    def pri(self, ide, rptr):
        print(ide)
        for i in rptr: print(i)

    def closeme(self):
        CLIENT.close()

    def sendfile(self):
        self.msg = '' # this is the message to go in the message line if anything fails
        scenario_names = self.scenarios.keys()
        try: # step 1 check the connection
            self.connectme()
        except:
            self.msg += 'Failed to find server.\n'
        else:
            self.commandme('rm /home/working/results/*')
            ftp_client=CLIENT.open_sftp()
            filestogo = glob.glob(self.K.outdir + '*')
            for f in filestogo:
                filenm = f.split('\\')[-1]
                try:
                    ftp_client.put(self.K.outdir + filenm, RESULTSDIR + filenm)
                except:
                    self.msg += 'Failed to send' + filenm + '\n'
            # done processing the files.  Close the connection.
            ftp_client.close()
            self.closeme()
        if self.msg == '':
            self.parent.frame.SetStatusText('Files sent')
        else:
            self.parent.frame.SetStatusText(self.msg)

    def getfile(self):
        try:
            self.connectme()
            ftp_client=CLIENT.open_sftp()
            fs = ftp_client.listdir(RESULTSDIR)
            for f in fs:
                ftp_client.get(RESULTSDIR + f, K.outdir +  f)
            ftp_client.close()
            self.closeme()
            self.parent.frame.SetStatusText('Files retrieved')
        except:
            self.parent.frame.msg('Failed at get file')

    def worker(self):
        try:
            self.connectme()
            self.commandme('./mp9starter.sh')
            self.closeme()
            self.parent.frame.SetStatusText('Process started on Microsoft Azure')
        except :
            self.parent.frame.msg('Worker failed')

    def getupdate(self):
        self.connectme()
        stdin, stdout, stderr = CLIENT.exec_command('tail ' + RESULTSDIR + 'output.txt')
        output = stdout.readlines()
        self.closeme()
        if len(output) > 0:
            lastline = output[-1]
        else:
            lastline = 'Running...'
        return lastline

    def checker(self):
        start = True
        while start:
            time.sleep(5)
            status = self.getupdate()
            self.parent.frame.SetStatusText(status)
            if status == 'FINISHED\n':
                start = False
        self.getfile()
        self.parent.frame.SetStatusText('Files retrieved and in output directory')


class Paars(wx.Panel):
    def __init__(self, parent, frame):
        wx.Panel.__init__(self, parent)
        self.parent = parent
        # Create the menubar and menu
        self.frame = frame 
        #empty data
        self.K = K
        self.K.stdfn, self.K.outdir = None, None
        self.status  = 'ok'
        self.census =  'census.json'
        self.scenario =  'scenario.json'
        # Now create the Panel to put the other controls on.
        panel = wx.Panel(self)
        lin = wx.StaticLine(panel,-1,style=wx.LI_HORIZONTAL)
        txt = wx.StaticText(panel, -1, " ")
        txta = wx.StaticText(panel, -1, " ")
        text = wx.StaticText(panel, -1, "Merit Insurance - Monte Carlo portfolio simulation")
        text.SetFont(wx.Font(14, wx.SWISS, wx.NORMAL, wx.BOLD))
        text.SetSize(text.GetBestSize())
        #self.state,self.color = getamastate()
        #self.statebtn = wx.ToggleButton(panel,-1,self.state)
        #self.statebtn.SetBackgroundColour(self.color)       
        #self.Bind(wx.EVT_TOGGLEBUTTON,self.switchamazon,self.statebtn)
        #self.statetext= wx.StaticText(panel, -1, "Click to switch Amazon on or off.  Currently...")
        #Now the filename and output directory
        self.rftext = wx.StaticText(panel, -1, "RollForward Spreadsheet "+ROLLFORWARD)
        self.rftext.SetFont(wx.Font(8, wx.SWISS, wx.NORMAL, wx.BOLD))
        self.sstext = wx.StaticText(panel, -1, "Input Spreadsheet")
        self.sstext.SetFont(wx.Font(8, wx.SWISS, wx.NORMAL, wx.BOLD))
        self.dirtext = wx.StaticText(panel, -1, "Output Directory")
        self.dirtext.SetFont(wx.Font(8, wx.SWISS, wx.NORMAL, wx.BOLD))
        # The buttons
        infil = wx.Button(panel, -1, "1. Choose the spreadsheet", (50,50))
        self.Bind(wx.EVT_BUTTON, self.choosefile, infil)
        outdir = wx.Button(panel, -1, "2. Choose the output directory", (50,50))
        self.Bind(wx.EVT_BUTTON, self.choosedirectory, outdir)
        self.processbtn = wx.Button(panel, -1, "3. Process...")
        self.Bind(wx.EVT_BUTTON, self.process, self.processbtn)
        rslt = wx.Button(panel, -1, "4. Show Summary Results...")
        self.Bind(wx.EVT_BUTTON, self.xldisplay, rslt)
        # Use a sizer to layout the controls
        sizer = wx.BoxSizer(wx.VERTICAL)
        sizer.Add(text, 0, wx.ALL, 10)
        sizer.Add(self.rftext, 0, wx.ALL, 10)
        sizer.Add(self.sstext, 0, wx.ALL, 10)
        sizer.Add(self.dirtext, 0, wx.ALL, 10)
        sizer.Add(txt)
        #sizer.Add(self.statetext, 0, wx.LEFT, 10)
        #sizer.Add(self.statebtn, 0, wx.ALL, 10)
        sizer.Add(txta)
        sizer.Add(infil, 0, wx.ALL, 10)
        sizer.Add(outdir, 0, wx.ALL, 10)
        sizer.Add(self.processbtn, 0, wx.ALL, 10)
        sizer.Add(rslt, 0, wx.ALL, 10)
        panel.SetSizer(sizer)
        panel.Layout()
        # And also use a sizer to manage the size of the panel such
        # that it fills the frame
        self.frame.SetStatusText('Remote process waiting')
        sizer = wx.BoxSizer()
        sizer.Add(panel, 1, wx.EXPAND)
        self.SetSizer(sizer)
        self.Fit()

    def setstate(self):
        self.state,self.color = getamastate()
        #print(self.state)
        self.statebtn.SetBackgroundColour(self.color)
        self.statebtn.SetLabel(self.state)

    def choosefile(self, evt):
        dlg = wx.FileDialog(
            self, message="Choose a file",
            defaultDir=os.getcwd(),
            defaultFile="",
            wildcard=WILDCARD,
            style=wx.FD_OPEN | 
                  wx.FD_CHANGE_DIR | wx.FD_FILE_MUST_EXIST |
                  wx.FD_PREVIEW
            )
        if dlg.ShowModal() == wx.ID_OK:
            path = dlg.GetPath()
            self.K.stdfn = path
            #self.K.stdfile = self.K.stdfn.split('\\')[-1]
            self.sstext.SetLabel("Input Spreadsheet: " + self.K.stdfn)
        dlg.Destroy()
        self.frame.SetStatusText('Input spreadsheet chosen.  Step 1 complete.')

    def do_spreadsheet(self):
        if self.K.stdfn and self.K.outdir:
            self.dfsc = pd.read_excel(self.K.stdfn,'Scenarios',index_col=0)
            self.dfsc.columns = [i.replace(' ','') for i in self.dfsc.columns]
            for col in self.dfsc.columns: # copy the rollforward spreadsheet for each run
                rollforward = col + '.xlsx'
                shutil.copy(ROLLFORWARD, self.K.outdir + rollforward)
                #print(ROLLFORWARD, self.K.outdir + rollforward)
            def_f = open(PCDIRECTORY + 'defaults.json')
            defaults = json.load(def_f)
            for k,v in defaults.items():
                if k not in self.dfsc.index: # if default value is already there, do not override
                    self.dfsc.loc[k] = v
            self.scenario_dick = self.dfsc.to_dict()
            self.dfsc.to_json(self.K.outdir + '\\' + self.scenario)
            self.frame.SetStatusText('Spreadsheet processed')
        else:
            self.frame.msg('You have not selected a spreadsheet and output directory yet')

    def xldisplay(self, evt):
        a = SummaryResult(self).ShowModal()

    def choosedirectory(self, evt):
        dlg = wx.DirDialog(self, "Choose a directory:",
                          style=wx.DD_DEFAULT_STYLE
                           )
        if dlg.ShowModal() == wx.ID_OK:
            self.K.outdir = dlg.GetPath() + '\\'
            self.dirtext.SetLabel("Output Directory: " + self.K.outdir)
        dlg.Destroy()
        self.frame.SetStatusText('Directory chosen,  Step 2 complete.')

    def process(self, evt):
        self.processbtn.Disable()
        self.do_spreadsheet()
        #self.moverollforward()
        montecarlo = Monty(self,self.scenario_dick,self.K)


class Blendo(wx.Panel):
    def __init__(self,parent,frame):
        wx.Panel.__init__(self,parent)
        self.parent = parent # refers back to the notebook
        self.K = K
        self.frame = frame
        self.K.outdir = '...'
        self.setuppath = '...' 
        self.vs = wx.StaticText(self, -1,'',(20,170))
        self.makesttext()
        # radio buttons to do Conning style or AAA style (annual return or cumulative value)
        lblList = ['Conning','AAA']
        self.filetype = lblList[0]
        self.rbox = wx.RadioBox(self, label = 'Who supplied the source files?', pos = (80,10), choices = lblList,
           majorDimension = 1, style = wx.RA_SPECIFY_ROWS) 
        self.rbox.Bind(wx.EVT_RADIOBOX,self.onRadioBox) 
        # end radio buttons 
        self.dirin = wx.Button(self, -1, "1. Choose the portfolio setup spreadsheet", (50,50))
        self.Bind(wx.EVT_BUTTON, self.choosesetupfile, self.dirin)
        self.dirout = wx.Button(self, -1, "2. Choose the output directory", (50,50))
        self.Bind(wx.EVT_BUTTON, self.chooseoutdirectory, self.dirout)
        self.Process = wx.Button(self, -1, "3. Process...")
        self.Bind(wx.EVT_BUTTON, self.process, self.Process)
        # Use a sizer to layout the controls, stacked vertically 10 pix border
        self.sizer = wx.BoxSizer(wx.VERTICAL)
        self.sizer.Add(self.vs, 0, wx.ALL, 10)
        self.sizer.Add(self.dirin, 0, wx.ALL, 10)
        self.sizer.Add(self.dirout, 0, wx.ALL, 10)
        self.sizer.Add(self.rbox, 0, wx.ALL, 10)
        self.sizer.Add(self.Process, 0,wx.ALL, 10)
        self.SetSizer(self.sizer)
        self.Layout()

    def makesttext(self):
        self.status = 'Portfolio setup files in ' + self.setuppath +'\nOutput csv file in ' + self.K.outdir
        self.vs.SetLabel(self.status)
   
    def onRadioBox(self,e): 
        self.filetype = self.rbox.GetStringSelection()

    def process(self,evt):
        df = pd.read_excel(self.setuppath,index_col=0,skiprows=1)
        self.filetype = self.rbox.GetStringSelection()
        if self.filetype == 'Conning':
            self.K.blendir = self.K.INDIR + 'ConningDec2020/'
        else:
            self.K.blendir = self.K.INDIR +  'AAACSV/'
        fail=False
        try:
            if self.filetype == 'Conning':
                trial = [
                         ['BondsGovtIntermediate',df.loc['BondsGovtIntermediate'].Weight],
                         ['BondsCorpLongInv',df.loc['BondsCorpLongInv'].Weight],
                         ['MoneyMarket',df.loc['MoneyMarket'].Weight],
                         ['BondsGovtShort',df.loc['BondsGovtShort'].Weight],
                         ['BondsGovtLong',df.loc['BondsGovtLong'].Weight],
                         ['BondsCorpShortInv',df.loc['BondsCorpShortInv'].Weight],
                         ['BondsCorpHighYield',df.loc['BondsCorpHighYield'].Weight],
                         ['BondsCorpIntermediateInv',df.loc['BondsCorpIntermediateInv'].Weight],
                         ['EquityUSLargeCap',df.loc['EquityUSLargeCap'].Weight],
                         ['EquityUSSmallCap',df.loc['EquityUSSmallCap'].Weight],
                         ['EquityUSAggressive',df.loc['EquityUSAggressive'].Weight],
                         ['EquityUSMidcap',df.loc['EquityUSMidcap'].Weight],
                         ['EquityInternationalAggressive',df.loc['EquityInternationalAggressive'].Weight],
                         ['EquityInternationalDiversified',df.loc['EquityInternationalDiversified'].Weight]]
            else:
                trial = [['US',df.loc['US'].Weight],
                         ['INT',df.loc['INT'].Weight],
                         ['SMALL',df.loc['SMALL'].Weight],
                         ['AGGR',df.loc['AGGR'].Weight],
                         ['MONEY',df.loc['MONEY'].Weight],
                         ['INTGOV',df.loc['INTGOV'].Weight],
                         ['LTCORP',df.loc['LTCORP'].Weight]]
        except:
            msg = 'The portfolio setup file is not correctly formatted.  Please review.'
            fail=True
        else:
            test = round(sum([i[1] for i in trial]),6)
            error = df.loc['St Dev Error'].Weight
            if test == 1.0:
                outputtab = [[i[0],"{:.2f}%".format(100*i[1])] for i in trial]
                msg= '<html><body><h1>Blending File Data</h1>'
                msg += tabulate(outputtab,['Asset','File','Weight %'],tablefmt='html')
                msg += '<p>Sum of weights = ' + "{:.2f}%".format(100*test) + '</p>'
                msg += '<p>Error = ' + "{:.2f}%".format(100*error) + '</p>'
                outfilenm = self.setuppath.split('\\')[-1].replace(' ','')
                outfilenm = outfilenm.replace('.xlsx','')
                outfilenm = outfilenm.replace('.xls','')
                outfilenm = outfilenm.upper().replace(' ','')
                outfilenm += 'ESG.csv' # PJE note sure if this will break something!
                self.K.csvfiles.append(outfilenm)
                msg += '<p>Output file is ' + outfilenm + '.csv</p>'
                msg += '<p>Close this window to process the file.</p>'
                msg +='</body></html>'
                htmlmsg("Blending file is ok...",msg)
                blendo(self, trial, self.K.blendir, self.K.outdir, outfilenm, error, self.filetype)
                self.frame.msg(outfilenm + ' created')
            else:
                msg = 'Portfolio totals do not equal 100% in weights.  Please fix and try again.'
                fail = True
        if fail:
            self.frame.msg(msg)

    def choosesetupfile(self, evt):
        dlg = wx.FileDialog(self, "Choose the input spreadsheet with the blends and error calc:",
                          style=wx.DD_DEFAULT_STYLE,
                          defaultDir=os.getcwd(),
                           )
        if dlg.ShowModal() == wx.ID_OK:
            self.setuppath = dlg.GetPath()
            self.makesttext()
        dlg.Destroy()

    def chooseoutdirectory(self, evt):
        dlg = wx.DirDialog(self, "Choose a directory for the output files:",
                          style=wx.DD_DEFAULT_STYLE
                           )
        if dlg.ShowModal() == wx.ID_OK:
            self.K.outdir = dlg.GetPath() + '\\'
            self.makesttext()
        dlg.Destroy()

class Easy(wx.Panel):
    def __init__(self,parent,frame):
        wx.Panel.__init__(self,parent)
        self.parent = parent # refers back to the notebook
        self.aaakeys = set(['Portfolio', 'date_range', 'Cash', 'SandP_500_Index', 'MSCI_EAFE_USD', 
                'BBgBarc_US_Treasury_3_5_Yr_TR_USD', 'Bbgbarc_US_Corporate_7_10_years_TR_USD', 
                'Russell_2000_TR_USD', 'AAAAggr', 'Style_R_Squared', 'Predicted_R_Squared'])
        self.connkeys = set(['Portfolio', 'date_range', 'Cash', 'ICE_BofA_5_10_Year_US_Treasury_Index', 
                'ICE_BofA_10plus_Year_US_Corporate_Index', 'ICE_BofA_1_5_Year_US_Treasury_Index', 
                'ICE_BofA_10plus_Year_US_Treasury_Index', 'ICE_BofA_1_5_Year_US_Corporate_Index', 
                'ICE_BofA_US_High_Yield_Index', 'ICE_BofA_5_10_Year_US_Corporate_Index', 
                'SandP_500_Index', 'R2000', 'NASDAQ_Composite_TR_USD', 'Russell_Mid_Cap_Index', 
                'MSCI_EM_Emerging_Markets_USD', 'MSCI_EAFE', 'Style_R_Squared', 'Predicted_R_Squared'])
        self.K = K
        self.frame = frame
        self.K.outdir = 'c:/PAARS/Portfolios/easy/'
        self.vs = wx.StaticText(self, -1,'',(20,170))
        self.makesttext()
        self.Dirout = wx.Button(self, -1, "1. Choose the output directory", (50,50))
        self.Bind(wx.EVT_BUTTON, self.chooseoutdirectory, self.Dirout)
        self.AAAPaste = wx.Button(self, -1, "2. Paste AAA Style MPI Data")
        self.Bind(wx.EVT_BUTTON, self.aaapaste, self.AAAPaste)
        # PJE leave this in - we might want Conning assets in future
        #self.ConningPaste = wx.Button(self, -1, "3. Paste Conning Style MPI Data")
        #self.Bind(wx.EVT_BUTTON, self.conningpaste, self.ConningPaste)
        self.Results = wx.Button(self, -1, "3. Show Results and Process Powerpoint")
        self.Bind(wx.EVT_BUTTON, self.results, self.Results)

        # Use a sizer to layout the controls, stacked vertically 10 pix border
        self.sizer = wx.BoxSizer(wx.VERTICAL)
        self.sizer.Add(self.vs, 0, wx.ALL, 10)
        self.sizer.Add(self.Dirout, 0, wx.ALL, 10)
        self.sizer.Add(self.AAAPaste, 0,wx.ALL, 10)
        #self.sizer.Add(self.ConningPaste, 0,wx.ALL, 10)
        self.sizer.Add(self.Results, 0,wx.ALL, 10)
        self.SetSizer(self.sizer)
        self.Layout()

    def makesttext(self):
        self.status = 'Output files in ' + self.K.outdir
        self.vs.SetLabel(self.status)

    def aaapaste(self, evt):
        self.error_msg, self.ass_weights, self.pftype = self.processClipboard()
        if self.pftype != 'AAA': self.error_msg = 'Please paste in AAA style assets'
        self.processblend()

    """ PJE leave this in as we might want to do Conning assets in the future
    def conningpaste(self, evt):
        self.error_msg, self.ass_weights, self.pftype = self.processClipboard()
        if self.pftype != 'Conning': self.error_msg = 'Please paste in Conning style assets'
        self.processblend()
    """

    def processblend(self):
        json_run = {} # this becomes scenario.json
        run_summary = {} # contains the weights and R_squared
        if not self.error_msg:
            for key,val in self.ass_weights.items():
                self.error_msg, fname, blendata = self.prepare_blend(self.error_msg, key, val, self.pftype)
                self.stderr = standarderr(blendata, self.pftype)
                run_summary[key] = blendata
                if not self.error_msg: blendo(self, blendata, self.K.blendir, self.K.outdir, fname, self.stderr, self.pftype) # puts ESG file in outdir
                if not self.error_msg: json_run[key] = self.json_details(key,fname)

        if not self.error_msg: # we have all the individual runs prepared
            self.frame.SetStatusText('Files blended ok.')
            with open(self.K.outdir + 'scenario.json','w') as jsfile:
                json.dump(json_run,jsfile)
            with open(self.K.outdir + 'summary.json','w') as summaryfile:
                json.dump(run_summary,summaryfile) # may need this for html display...
            shutil.copy(PCDIRECTORY+'census.json',self.K.outdir+'census.json') # always the same census

        if self.error_msg:
            wx.MessageBox(self.error_msg,"Error")
        else:
            montecarlo = Monty(self,json_run,self.K) # start the Monte Carlo engine

    def results(self, evt):
        a = SummaryResult(self) #.Show() #Modal() # create the results table
        #a.Destroy()
        print('1')
        results_pic() # put the results table to a png file
        print('2')
        histo() # put the spectrumthisy*.P files to a histogram png
        print('3')
        outputfile = outputppt(self.K.outdir) # add a slide to the MPI power point or create a 1 slide ppt
        print('4')
        wx.MessageBox('Created summary powerpoint slide and put it in\n' + outputfile,"Message")
        print('5')


    def processClipboard(self):
        data = wx.TextDataObject()
        if wx.TheClipboard.Open():
            success = wx.TheClipboard.GetData(data)
            wx.TheClipboard.Close()
        if not success:
            failmsg = 'No data in clipboard'
        else:        
            failmsg = '' # this is success - a blank failmsg

        try:
            lines = data.GetText().split('\n')
        except:
            failmsg = 'Clipboard does not contain data'
        if not failmsg:
            try:
                heads = lines[1].split('\t')
                heads[1] = 'date_range'
                heads[0] = 'Portfolio'
            except:
                failmsg = 'Invalid clipboard, no MPI heading data'
        if not failmsg:
            delete_dict = {' ':'_','|':'_','-':'_',')':'','(':'','+':'plus','&':'and'}
            deletetable = str.maketrans(delete_dict)
            heads = [i.translate(deletetable) for i in heads]
            try:
                pftuple = namedtuple('portfolio',heads)
            except:
                failmsg = 'Invalid paste column headers.  Did you paste the right area?'
        if not failmsg:
            clipdata = OrderedDict()
            for line in lines[2:]: # may need a try/except here, see if it ever fails...
                linedata = line.split('\t')
                linedata = [flt(i) for i in linedata]
                try:
                    for i in [0,1]: linedata[i] = linedata[i].translate(deletetable)
                    pfdata = pftuple(*linedata)
                    clipdata[pfdata.Portfolio] = pfdata
                except:
                    pass # there is a harmless blank line at the end, which will cause an error

            if set(heads) == self.aaakeys:
                pftype = 'AAA'
                retkeys = self.aaakeys
            elif set(heads) == self.connkeys:
                pftype = 'Conning'
                retkeys = self.connkeys
            else:
                pftype = 'Unknown'
                retkeys = None
                failmsg = 'Unknown Portfolio keys ' + ' '.join(heads)
            lst = list(clipdata.keys())
            defaults = range(len(lst)) # default choice = everything
            dlg = wx.MultiChoiceDialog( self, # now set up a chooser, so you do not have to do it all
                               "Choose the portfolios",
                               pftype + "Choices...", lst)
            dlg.SetSelections(defaults)
            if (dlg.ShowModal() == wx.ID_OK):
                selections = dlg.GetSelections()
                keysneeded = [lst[x] for x in selections]
            else:
                keysneeded = None
            dlg.Destroy()
            result = dict((k, clipdata[k]) for k in keysneeded)
        else:
            result = None
            pftype = 'Not known'
        return failmsg, result, pftype

    def prepare_blend(self, err_msg, name, data, pftype):
        msg = ''
        self.filetype = pftype
        if self.filetype == 'Conning':
            self.K.blendir = self.K.INDIR + 'ConningDec2020/'
        else:
            self.K.blendir = self.K.INDIR +  'AAACSV/'
        fail=False
        try:
            if self.filetype == 'Conning':
                output = [
                         ('MoneyMarket',data.Cash),
                         ('BondsCorpHighYield',data.ICE_BofA_US_High_Yield_Index),
                         ('BondsCorpIntermediateInv',data.ICE_BofA_5_10_Year_US_Corporate_Index),
                         ('BondsCorpLongInv',data.ICE_BofA_10plus_Year_US_Corporate_Index),
                         ('BondsCorpShortInv',data.ICE_BofA_1_5_Year_US_Corporate_Index),
                         ('BondsGovtIntermediate',data.ICE_BofA_5_10_Year_US_Treasury_Index),
                         ('BondsGovtLong',data.ICE_BofA_10plus_Year_US_Treasury_Index),
                         ('BondsGovtShort',data.ICE_BofA_1_5_Year_US_Treasury_Index),
                         ('EquityInternationalAggressive',data.MSCI_EM_Emerging_Markets_USD),
                         ('EquityInternationalDiversified',data.MSCI_EAFE),
                         ('EquityUSAggressive',data.NASDAQ_Composite_TR_USD),
                         ('EquityUSLargeCap',data.SandP_500_Index),
                         ('EquityUSMidcap',data.Russell_Mid_Cap_Index),
                         ('EquityUSSmallCap',data.R2000)] #)
            else:
                output = [ 
                         ('US',data.SandP_500_Index),
                         ('INT',data.MSCI_EAFE_USD),
                         ('SMALL',data.Russell_2000_TR_USD),
                         ('AGGR',data.AAAAggr),
                         ('MONEY',data.Cash),
                         ('INTGOV',data.BBgBarc_US_Treasury_3_5_Yr_TR_USD),
                         ('LTCORP',data.Bbgbarc_US_Corporate_7_10_years_TR_USD)] #)
        except:
            msg += 'The portfolio setup file is not correctly formatted.  Please review.'
        output.append(('total', sum([i[1] for i in output])))
        output.append(('Predicted_R_Squared', data.Predicted_R_Squared))
        output2 = [(i[0],round(i[1]/100,6)) for i in output]
        outfilenm = name.upper().replace(' ','') 
        outcsv = outfilenm + 'ESG.csv'
        rollforward = outfilenm + '.xlsx'
        shutil.copy(ROLLFORWARD, self.K.outdir + rollforward)
        return msg, outcsv, output2

    def json_details(self,name,fname):
        return {"Description": name,
                "CsvFile": fname,
                "CensusFile": None,
                "Fund":100000000.0,
                "Population": 1000,
                "Runs": 1000,
                "AnnCum": "CUM",
                "Income": .05,
                "Premium": .0055,
                "WMFee": .01,
                "DiscountRate": .03,
                "MortSpreadsheet": "IAM20122581_2582.xlsx",
                "LapseUtilization":"Utilization2020_1_26.xlsx",
                "Stochastic": False,
                "Prudent": True,
                "Debug":False,
                "YearOneIncome": 1.0}

    def chooseoutdirectory(self, evt):
        dlg = wx.DirDialog(self, "Choose a directory for the output files:",
                          style=wx.DD_DEFAULT_STYLE
                           )
        if dlg.ShowModal() == wx.ID_OK:
            self.K.outdir = dlg.GetPath() + '\\'
            self.makesttext()
        dlg.Destroy()

class MainFrame(wx.Frame):
    def __init__(self):
        wx.Frame.__init__(self,None,title="MeritInsurance Monte Carlo Simulator",size=(1000,700))
        p = wx.Panel(self)
        self.Bind(wx.EVT_CLOSE, self.OnClose) # event close
        nb = wx.Notebook(p)
        self.CreateStatusBar()
        paarstab = Paars(nb,self)
        blendtab = Blendo(nb,self)
        easytab = Easy(nb,self)
        menuBar = wx.MenuBar()
        menu = wx.Menu()
        menu.Append(101, "&Help\tAlt-H","This will show the instructions")
        menu.Append(wx.ID_EXIT, "E&xit\tAlt-X", "Exit")
        self.Bind(wx.EVT_MENU, self.Help, id=101)
        self.Bind(wx.EVT_MENU, self.Exit, id=wx.ID_EXIT)
        menuBar.Append(menu,"&Help")
        self.SetMenuBar(menuBar)
        nb.AddPage(paarstab,"Run Paars")
        nb.AddPage(blendtab,"Create Scenarios from blended asset factors")
        nb.AddPage(easytab,"QuickRun: one MPI table")
        sizer = wx.BoxSizer()
        sizer.Add(nb,1,wx.EXPAND)
        p.SetSizer(sizer)
        p.Layout()
        p.Fit()

    def OnClose(self, event):
        #self.state,self.color = getamastate()
        #if event.CanVeto() and self.state in ['Running','Pending']:
        #if self.state in ['running','pending']:
        #    if wx.MessageBox("The environment has not been stopped... continue closing?",
        #        "Please confirm",
        #        wx.ICON_QUESTION | wx.YES_NO) != wx.YES:
        #        event.Veto()
        #        return
        self.Destroy()  # you may also do:  event.Skip()
			# since the default event handler does call Destroy(), too


    def msg(self, msg):
        dlg = wx.MessageDialog(self, msg,
                                     'Message',
                                      wx.OK | wx.ICON_INFORMATION
                                      )
        dlg.ShowModal()
        dlg.Destroy()

    def Help(self, evt):
        text = '''There are four steps...\n
First, set up the spreadsheet with the scenarios and census.  
\tYou can use the spreadsheet Standard.xlsx as a template
\tThen choose the file you have set up.\n
Second, choose an output directory.\n
Third, the Process button will check the spreadsheet,
\tsend it to Amazon for processing, and
\tsend any output back to this computer.\n
Fourth you can review the graphs and spreadsheet
and run the dashboard with the files in the Output directory.
        '''
        self.msg(text)
        #dlg = wx.MessageDialog(self, text,
        #                             'Message',
        #                              wx.OK | wx.ICON_INFORMATION
        #                              )
        #dlg.ShowModal()
        #dlg.Destroy()

    def Exit(self, evt):
        """Event handler for the button click."""
        self.Close()

if __name__ == "__main__":
    app = wx.App()
    MainFrame().Show()
    app.MainLoop()

